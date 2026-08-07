"""Readers for the formats a person calls "a document".

PDF, Word, Excel, PowerPoint, OpenDocument, EPUB, RTF, CSV, and anything that
is already text. The PDF reader is the one with an opinion worth reading: a PDF
with no text layer is a photograph of a page, and pretending otherwise — by
returning an empty string and calling it a success — is how a scanned invoice
becomes "the document appears to be blank".

Optional dependencies are imported inside the function that needs them, with
the install line in the warning. Annona is installed on machines that will
never open a .pptx, and making every one of them carry python-pptx to read a
text file is the kind of default that gets a tool removed from the image.
"""

from __future__ import annotations

import csv
import re
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

from loguru import logger

from runner.tools.extractors.images import ocr_image
from runner.tools.extractors.types import Extraction, MediaRef, ReadOptions, missing_dependency

__all__ = [
    "read_csv",
    "read_epub",
    "read_excel",
    "read_odf",
    "read_pdf",
    "read_powerpoint",
    "read_rtf",
    "read_text",
    "read_word",
    "render_first_page",
    "strip_html",
    "tidy",
]

_HTML_TAG = re.compile(r"<[^>]+>")
_RTF_CONTROL = re.compile(r"\\[a-zA-Z]+-?\d* ?|[{}]|\\\n")
_BLANK_RUN = re.compile(r"\n{3,}")


def tidy(text: str) -> str:
    """Collapse the blank-line runs that every converter leaves behind."""
    return _BLANK_RUN.sub("\n\n", text).strip()


def strip_html(markup: str) -> str:
    """Tags out, words kept. Crude on purpose — a mail client's HTML is furniture."""
    return tidy(_HTML_TAG.sub(" ", markup))


# ── PDF ───────────────────────────────────────────────────────────────────────


def read_pdf(path: Path, opts: ReadOptions) -> Extraction:
    """Text layer first, OCR second, and the truth about which one answered."""
    pages: list[str] = []
    metadata: dict = {"format": "pdf"}

    try:
        import pdfplumber  # noqa: PLC0415 — heavy, and only PDFs need it

        with pdfplumber.open(str(path)) as pdf:
            metadata["pages"] = len(pdf.pages)
            metadata["info"] = {k: str(v) for k, v in (pdf.metadata or {}).items()}
            for number, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"--- Page {number} ---\n{text}")
    except ImportError:
        from pypdf import PdfReader  # noqa: PLC0415 — the fallback path

        reader = PdfReader(str(path))
        metadata["pages"] = len(reader.pages)
        for number, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- Page {number} ---\n{text}")

    extraction = Extraction(
        format="pdf",
        text="\n\n".join(pages),
        metadata=metadata,
        # Offered whether or not there was a text layer: a substrate that can
        # read documents natively sees the tables and the stamps, which no text
        # extractor recovers.
        media=(MediaRef(path=path, media_type="pdf", label=path.name),),
    )

    if extraction.text.strip():
        return extraction

    ocr_text, ocr_note = _ocr_pdf(path, opts)
    if ocr_text:
        return Extraction(
            format="pdf",
            text=ocr_text,
            metadata={**metadata, "text_layer": False, "ocr": True},
            media=extraction.media,
            warnings=("this PDF has no text layer; the text above came from OCR",),
        )

    return extraction.warn(
        "this PDF has no text layer — it is a scan. "
        + (ocr_note or "Nothing was read from it as text.")
    )


def _ocr_pdf(path: Path, opts: ReadOptions) -> tuple[str, str]:
    """Rasterise and OCR, if this machine happens to be able to.

    Returns the text and, when there is none, the reason. Both halves matter:
    an empty result with no explanation is indistinguishable from an empty page.
    """
    if opts.ocr == "never":
        return "", "OCR was disabled for this read."

    try:
        import pypdfium2  # noqa: PLC0415 — optional rasteriser
    except ImportError:
        return "", missing_dependency("pypdfium2", "OCR of a scanned PDF")

    out: list[str] = []
    derived = opts.derived_dir(path)
    document = pypdfium2.PdfDocument(str(path))
    try:
        for index in range(min(len(document), 20)):
            image_path = derived / f"{path.stem}-p{index + 1}.png"
            document[index].render(scale=2).to_pil().save(image_path)
            text, note = ocr_image(image_path, opts)
            if not text and note:
                return "", note
            if text.strip():
                out.append(f"--- Page {index + 1} (OCR) ---\n{text}")
    finally:
        document.close()

    return "\n\n".join(out), ""


def render_first_page(path: Path, target: Path) -> Path | None:
    """The first page of a PDF as an image, or ``None`` if this build cannot.

    Used for the thumbnail in the window. A PDF is the one document format where
    the cover *is* the identity — an operator recognises the invoice by its
    letterhead long before they read the number.
    """
    try:
        import pypdfium2  # noqa: PLC0415 — optional rasteriser
    except ImportError:
        return None

    try:
        document = pypdfium2.PdfDocument(str(path))
    except Exception as exc:  # noqa: BLE001 — a broken PDF has no cover, that is all
        logger.debug(f"could not open {path} for a preview: {exc}")
        return None

    try:
        if len(document) == 0:
            return None
        target.parent.mkdir(parents=True, exist_ok=True)
        document[0].render(scale=1.2).to_pil().convert("RGB").save(target)
        return target
    finally:
        document.close()


# ── Office ────────────────────────────────────────────────────────────────────


def read_word(path: Path, opts: ReadOptions) -> Extraction:
    """Paragraphs and tables. A contract's obligations are usually in the tables."""
    try:
        from docx import Document  # noqa: PLC0415 — optional
    except ImportError:
        return Extraction(
            format="docx", warnings=(missing_dependency("python-docx", "reading Word documents"),)
        )

    document = Document(str(path))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    tables = [
        "\n".join(" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows)
        for table in document.tables
    ]

    text = "\n".join(paragraphs)
    if tables:
        text += "\n\n=== Tables ===\n" + "\n\n---\n".join(tables)

    return Extraction(
        format="docx",
        text=tidy(text),
        metadata={"format": "docx", "paragraphs": len(paragraphs), "tables": len(tables)},
    )


def read_excel(path: Path, opts: ReadOptions) -> Extraction:
    """Every sheet unless one was named. Values, not formulas."""
    try:
        import openpyxl  # noqa: PLC0415 — optional
    except ImportError:
        return Extraction(
            format="xlsx", warnings=(missing_dependency("openpyxl", "reading Excel workbooks"),)
        )

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sheets = workbook.sheetnames
    wanted = [opts.sheet_name] if opts.sheet_name in sheets else sheets

    blocks: list[str] = []
    for name in wanted:
        rows = [
            " | ".join("" if cell is None else str(cell) for cell in row)
            for row in workbook[name].iter_rows(values_only=True)
        ]
        rows = [row for row in rows if row.strip(" |")]
        if rows:
            blocks.append(f"=== Sheet: {name} ===\n" + "\n".join(rows))

    return Extraction(
        format="xlsx",
        text="\n\n".join(blocks),
        metadata={"format": "xlsx", "sheets": sheets, "sheets_read": wanted},
    )


def read_powerpoint(path: Path, opts: ReadOptions) -> Extraction:
    """Slide text and speaker notes, in slide order."""
    try:
        from pptx import Presentation  # noqa: PLC0415 — optional
    except ImportError:
        return Extraction(
            format="pptx", warnings=(missing_dependency("python-pptx", "reading presentations"),)
        )

    deck = Presentation(str(path))
    blocks: list[str] = []
    for number, slide in enumerate(deck.slides, 1):
        lines = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append(f"[notes] {slide.notes_slide.notes_text_frame.text.strip()}")
        if lines:
            blocks.append(f"--- Slide {number} ---\n" + "\n".join(lines))

    return Extraction(
        format="pptx",
        text="\n\n".join(blocks),
        metadata={"format": "pptx", "slides": len(deck.slides)},
    )


def read_odf(path: Path, opts: ReadOptions) -> Extraction:
    """OpenDocument text, spreadsheet or presentation, read as what it is: a zip.

    No dependency, because an ODF file is XML in a zip container and the text
    is in one element. Formatting is lost; the words are not.
    """
    kind = {".odt": "odt", ".ods": "ods", ".odp": "odp"}.get(path.suffix.lower(), "odf")
    with zipfile.ZipFile(path) as archive:
        try:
            payload = archive.read("content.xml")
        except KeyError:
            return Extraction(format=kind, warnings=("no content.xml inside the document",))

    root = ET.fromstring(payload)
    text = "\n".join(part.strip() for part in root.itertext() if part.strip())
    return Extraction(format=kind, text=tidy(text), metadata={"format": kind})


def read_epub(path: Path, opts: ReadOptions) -> Extraction:
    """An EPUB is a zip of XHTML. Read it in spine order, tags stripped."""
    chapters: list[str] = []
    with zipfile.ZipFile(path) as archive:
        names = [n for n in archive.namelist() if n.lower().endswith((".xhtml", ".html", ".htm"))]
        for name in sorted(names)[: opts.max_members]:
            body = archive.read(name).decode("utf-8", errors="replace")
            stripped = strip_html(body)
            if stripped:
                chapters.append(f"--- {name} ---\n{stripped}")

    return Extraction(
        format="epub",
        text="\n\n".join(chapters),
        metadata={"format": "epub", "documents": len(chapters)},
    )


# ── Plain formats ─────────────────────────────────────────────────────────────


def read_csv(path: Path, opts: ReadOptions) -> Extraction:
    rows: list[str] = []
    with open(path, newline="", encoding="utf-8", errors="replace") as handle:
        # Sniffing beats assuming a comma: European exports are semicolon-
        # separated far more often than not, and a wrong delimiter turns a
        # readable table into one column of noise.
        sample = handle.read(8192)
        handle.seek(0)
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except csv.Error:
            dialect = csv.excel
        for row in csv.reader(handle, dialect):
            rows.append(" | ".join(row))

    return Extraction(
        format="csv",
        text="\n".join(rows),
        metadata={
            "format": "csv",
            "rows": len(rows),
            "delimiter": getattr(dialect, "delimiter", ","),
        },
    )


def read_rtf(path: Path, opts: ReadOptions) -> Extraction:
    """RTF without a dependency: strip the control words, keep the words."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    text = _RTF_CONTROL.sub("", raw)
    return Extraction(format="rtf", text=tidy(text), metadata={"format": "rtf"})


def read_text(path: Path, opts: ReadOptions) -> Extraction:
    text = path.read_text(encoding="utf-8", errors="replace")
    fmt = path.suffix.lstrip(".").lower() or "txt"
    logger.debug(f"read {path} as text ({len(text)} chars)")
    return Extraction(
        format=fmt, text=text, metadata={"format": fmt, "lines": text.count("\n") + 1}
    )
